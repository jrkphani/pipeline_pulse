"""
O2R Export Service - Generate slides, decks, and reports
"""

import io
import tempfile
from typing import List, Dict, Any, Optional
from datetime import datetime, date
from pathlib import Path
import base64

# For PDF generation
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.piecharts import Pie
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# For PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

from app.models.o2r.opportunity import O2ROpportunity, OpportunityPhase, HealthSignalType
from app.models.o2r.health import HealthSignalEngine

class O2RExportService:
    """Service for exporting O2R data in various formats"""
    
    def __init__(self):
        self.health_engine = HealthSignalEngine()
    
    def generate_opportunity_slide(self, opportunity: O2ROpportunity) -> bytes:
        """Generate single opportunity slide as PDF"""
        
        if not REPORTLAB_AVAILABLE:
            raise ImportError("ReportLab is required for PDF generation")
        
        # Create PDF in memory
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=1  # Center
        )
        
        header_style = ParagraphStyle(
            'CustomHeader',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=10,
            textColor=colors.blue
        )
        
        # Title
        title = f"{opportunity.deal_name}"
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))
        
        # Basic info table
        basic_info = [
            ['Account', opportunity.account_name],
            ['Owner', opportunity.owner],
            ['Value', f"SGD {opportunity.sgd_amount:,.0f}"],
            ['Probability', f"{opportunity.probability}%"],
            ['Territory', opportunity.territory or 'N/A'],
            ['Service Type', opportunity.service_type or 'N/A'],
            ['Current Phase', opportunity.current_phase.replace('_', ' ').title()],
            ['Health Signal', self._format_health_signal(opportunity.health_signal)]
        ]
        
        basic_table = Table(basic_info, colWidths=[2*inch, 3*inch])
        basic_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.lightblue),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('BACKGROUND', (1, 0), (1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(basic_table)
        story.append(Spacer(1, 30))
        
        # Phase timeline
        story.append(Paragraph("Phase Timeline", header_style))
        phase_data = self._generate_phase_timeline_data(opportunity)
        
        phase_table = Table(phase_data, colWidths=[1.5*inch, 1.5*inch, 2*inch])
        phase_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.blue),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(phase_table)
        story.append(Spacer(1, 20))
        
        # Action items
        if opportunity.action_items:
            story.append(Paragraph("Action Items", header_style))
            for item in opportunity.action_items:
                story.append(Paragraph(f"• {item}", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Blockers
        if opportunity.blockers:
            story.append(Paragraph("Current Blockers", header_style))
            for blocker in opportunity.blockers:
                story.append(Paragraph(f"• {blocker}", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Comments
        if opportunity.comments:
            story.append(Paragraph("Comments", header_style))
            story.append(Paragraph(opportunity.comments, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        pdf_bytes = buffer.getvalue()
        buffer.close()
        
        return pdf_bytes
    
    def generate_territory_deck(self, territory: str, opportunities: List[O2ROpportunity]) -> bytes:
        """Generate PowerPoint deck for territory"""
        
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint generation")
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{territory} O2R Pipeline Review"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Summary slide
        self._add_territory_summary_slide(prs, territory, opportunities)
        
        # Individual opportunity slides
        for opp in opportunities:
            self._add_opportunity_slide(prs, opp)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_bytes = buffer.getvalue()
        buffer.close()
        
        return pptx_bytes
    
    def generate_weekly_review_report(self, review_data: Dict[str, Any]) -> bytes:
        """Generate weekly review report as PDF"""
        
        if not REPORTLAB_AVAILABLE:
            raise ImportError("ReportLab is required for PDF generation")
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'ReviewTitle',
            parent=styles['Heading1'],
            fontSize=28,
            spaceAfter=30,
            alignment=1,
            textColor=colors.blue
        )
        
        title = f"Weekly O2R Review - {review_data['week_of']}"
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 30))
        
        # Summary section
        summary = review_data['summary']
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        
        summary_data = [
            ['Total Opportunities', f"{summary['total_opportunities']}"],
            ['Total Value', f"SGD {summary['total_value_sgd']:,.0f}"],
            ['Updated This Week', f"{summary['updated_this_week']} ({summary.get('update_percentage', 0):.1f}%)"],
            ['Require Attention', f"{summary['attention_required']}"],
            ['Revenue Realized', f"SGD {summary['revenue_realized_this_week']:,.0f}"],
            ['New Opportunities', f"{summary['new_opportunities']}"]
        ]
        
        summary_table = Table(summary_data, colWidths=[2.5*inch, 2*inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.lightblue),
            ('TEXTCOLOR', (0, 0), (0, -1), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 11),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('BACKGROUND', (1, 0), (1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(summary_table)
        story.append(Spacer(1, 30))
        
        # Key highlights
        insights = review_data['insights']
        if insights['key_highlights']:
            story.append(Paragraph("Key Highlights", styles['Heading2']))
            for highlight in insights['key_highlights']:
                story.append(Paragraph(f"• {highlight}", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Concerns
        if insights['concerns']:
            story.append(Paragraph("Areas of Concern", styles['Heading2']))
            for concern in insights['concerns']:
                story.append(Paragraph(f"• {concern}", styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Action items
        action_items = review_data['action_items']
        if action_items:
            story.append(Paragraph("Action Items", styles['Heading2']))
            
            # Group by priority
            high_priority = [ai for ai in action_items if ai['priority'] == 'High']
            medium_priority = [ai for ai in action_items if ai['priority'] == 'Medium']
            
            if high_priority:
                story.append(Paragraph("High Priority", styles['Heading3']))
                for ai in high_priority:
                    story.append(Paragraph(f"• {ai['owner']}: {ai['description']}", styles['Normal']))
                story.append(Spacer(1, 10))
            
            if medium_priority:
                story.append(Paragraph("Medium Priority", styles['Heading3']))
                for ai in medium_priority:
                    story.append(Paragraph(f"• {ai['owner']}: {ai['description']}", styles['Normal']))
        
        # Build PDF
        doc.build(story)
        pdf_bytes = buffer.getvalue()
        buffer.close()
        
        return pdf_bytes
    
    def generate_csv_export(self, opportunities: List[O2ROpportunity]) -> str:
        """Generate CSV export of opportunities"""
        
        import csv
        import io
        
        output = io.StringIO()
        writer = csv.writer(output)
        
        # Headers
        headers = [
            'ID', 'Deal Name', 'Account', 'Owner', 'SGD Amount', 'Probability', 
            'Current Stage', 'Current Phase', 'Territory', 'Service Type', 
            'Funding Type', 'Health Signal', 'Health Reason',
            'Proposal Date', 'PO Date', 'Kickoff Date', 'Invoice Date', 
            'Payment Date', 'Revenue Date', 'Last Updated', 'Comments'
        ]
        writer.writerow(headers)
        
        # Data rows
        for opp in opportunities:
            row = [
                opp.id,
                opp.deal_name,
                opp.account_name,
                opp.owner,
                opp.sgd_amount,
                opp.probability,
                opp.current_stage,
                opp.current_phase,
                opp.territory or '',
                opp.service_type or '',
                opp.funding_type or '',
                opp.health_signal,
                opp.health_reason or '',
                opp.proposal_date.isoformat() if opp.proposal_date else '',
                opp.po_date.isoformat() if opp.po_date else '',
                opp.kickoff_date.isoformat() if opp.kickoff_date else '',
                opp.invoice_date.isoformat() if opp.invoice_date else '',
                opp.payment_date.isoformat() if opp.payment_date else '',
                opp.revenue_date.isoformat() if opp.revenue_date else '',
                opp.last_updated.isoformat(),
                opp.comments or ''
            ]
            writer.writerow(row)
        
        csv_content = output.getvalue()
        output.close()
        
        return csv_content
    
    def _format_health_signal(self, signal: HealthSignalType) -> str:
        """Format health signal for display"""
        signal_map = {
            HealthSignalType.GREEN: "🟢 On Track",
            HealthSignalType.YELLOW: "🟡 Attention",
            HealthSignalType.RED: "🔴 Critical",
            HealthSignalType.BLOCKED: "🚧 Blocked",
            HealthSignalType.NEEDS_UPDATE: "📋 Needs Update"
        }
        return signal_map.get(signal, signal)
    
    def _generate_phase_timeline_data(self, opportunity: O2ROpportunity) -> List[List[str]]:
        """Generate phase timeline data for tables"""
        
        phases = [
            ("Phase 1", "Opportunity to Proposal", opportunity.proposal_date),
            ("Phase 2", "Proposal to Commitment", opportunity.po_date),
            ("Phase 3", "Execution", opportunity.kickoff_date),
            ("Phase 4", "Revenue Realization", opportunity.revenue_date)
        ]
        
        data = [["Phase", "Description", "Key Date"]]
        
        for phase_name, description, key_date in phases:
            date_str = key_date.strftime("%Y-%m-%d") if key_date else "Pending"
            data.append([phase_name, description, date_str])
        
        return data
    
    def _add_territory_summary_slide(self, prs, territory: str, opportunities: List[O2ROpportunity]):
        """Add territory summary slide to presentation"""
        
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{territory} Summary"
        
        # Summary metrics
        total_value = sum(opp.sgd_amount for opp in opportunities)
        avg_deal_size = total_value / len(opportunities) if opportunities else 0
        health_summary = self.health_engine.get_health_summary_for_portfolio(opportunities)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        p = tf.paragraphs[0]
        p.text = f"Total Opportunities: {len(opportunities)}"
        
        p = tf.add_paragraph()
        p.text = f"Total Value: SGD {total_value:,.0f}"
        
        p = tf.add_paragraph()
        p.text = f"Average Deal Size: SGD {avg_deal_size:,.0f}"
        
        p = tf.add_paragraph()
        p.text = f"Health Summary: {health_summary['green']} Green, {health_summary['yellow']} Yellow, {health_summary['red']} Red"
    
    def _add_opportunity_slide(self, prs, opportunity: O2ROpportunity):
        """Add individual opportunity slide to presentation"""
        
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = opportunity.deal_name
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Basic info
        p = tf.paragraphs[0]
        p.text = f"Account: {opportunity.account_name}"
        
        p = tf.add_paragraph()
        p.text = f"Value: SGD {opportunity.sgd_amount:,.0f}"
        
        p = tf.add_paragraph()
        p.text = f"Owner: {opportunity.owner}"
        
        p = tf.add_paragraph()
        p.text = f"Phase: {opportunity.current_phase.replace('_', ' ').title()}"
        
        p = tf.add_paragraph()
        p.text = f"Health: {self._format_health_signal(opportunity.health_signal)}"
        
        # Next milestone
        next_milestone = opportunity.get_next_milestone()
        if next_milestone:
            p = tf.add_paragraph()
            p.text = f"Next: {next_milestone.replace('_', ' ').title()}"
    
    def export_to_base64(self, data: bytes, filename: str) -> Dict[str, str]:
        """Convert bytes to base64 for API response"""
        base64_data = base64.b64encode(data).decode('utf-8')
        
        # Determine MIME type
        mime_type = "application/octet-stream"
        if filename.endswith('.pdf'):
            mime_type = "application/pdf"
        elif filename.endswith('.pptx'):
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif filename.endswith('.csv'):
            mime_type = "text/csv"
        
        return {
            'filename': filename,
            'mime_type': mime_type,
            'data': base64_data,
            'size': len(data)
        }
